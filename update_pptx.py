"""
Script to update HardwarePipeline_FINAL_v3.pptx based on SLIDE_REVIEW_NOTES.md corrections

Requirements:
    pip install python-pptx

Usage:
    python update_pptx.py path/to/HardwarePipeline_FINAL_v3.pptx
"""

import sys
from pathlib import Path
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
except ImportError:
    print("Installing python-pptx...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

def update_pptx(pptx_path: Path):
    """Update the PPTX with all corrections from SLIDE_REVIEW_NOTES.md"""

    print(f"Loading {pptx_path}...")
    prs = Presentation(str(pptx_path))

    print(f"Slides before: {len(prs.slides)}")

    # Track changes
    changes = []

    # === SLIDE 1: Title / Phase Overview ===
    print("\n=== SLIDE 1: Phase Overview ===")

    # Change "8 Phases" to "11 Phases"
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                if "8 Phases" in text or "8 phases" in text or "8 phase" in text.lower():
                    new_text = text.replace("8 Phases", "11 Phases").replace("8 phases", "11 phases")
                    shape.text = new_text
                    changes.append(f"Slide {slide.slide_id}: '8 Phases' => '11 Phases'")
                    print(f"  [OK] Updated: {text[:50]}... => {new_text[:50]}...")

    # === SLIDE 2: Problem Statement ===
    print("\n=== SLIDE 2: Problem Statement ===")

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text

                # A3: "from requirements to code review" => "from requirements to production-ready firmware"
                if "from requirements to code review" in text.lower():
                    new_text = text.replace("from requirements to code review", "from requirements to production-ready firmware")
                    shape.text = new_text
                    changes.append(f"Slide {slide.slide_id}: Rephrased bullet 01")
                    print(f"  [OK] Updated: '{text}'")

                # A6: Replace FCC with MIL-STD, add IEEE 830, 1016
                if "FCC" in text and "MISRA-C" in text:
                    new_text = text.replace("FCC, MISRA-C", "MIL-STD-461/810, MISRA-C")
                    if "IEEE 29148" in text:
                        new_text = new_text.replace("IEEE 29148", "IEEE 29148, IEEE 830, IEEE 1016")
                    elif "IEEE" in text:
                        new_text = new_text.replace("IEEE", "IEEE 830, IEEE 1016, IEEE")
                    else:
                        new_text = new_text.replace("MIL-STD-461/810, MISRA-C",
                                                          "MIL-STD-461/810, IEEE 29148, IEEE 830, IEEE 1016, MISRA-C")
                    shape.text = new_text
                    changes.append(f"Slide {slide.slide_id}: Updated standards list")
                    print(f"  [OK] Updated standards")

                # A4: Add footnote for INR42L calculation
                if "INR42L" in text or "42 Lakhs" in text:
                    # Check if footnote already exists
                    if "* Based on" not in text:
                        new_text = text + "\n\n* Based on MIL-grade PCB respin: fab + components + engineering time (~$50K / INR84)"
                        shape.text = new_text
                        changes.append(f"Slide {slide.slide_id}: Added INR42L footnote")
                        print(f"  [OK] Added footnote for INR42L")

    # === SLIDE 3: Proposed Solution ===
    print("\n=== SLIDE 3: Solution ===")

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text

                # A9: Remove "port 8000"
                if "port 8000" in text.lower():
                    new_text = text.replace(" - port 8000", "").replace(": 8000", "")
                    if "FASTAPI" in new_text:
                        new_text = new_text.replace("FASTAPI", "FASTAPI BACKEND · LOCAL REST API")
                    shape.text = new_text
                    changes.append(f"Slide {slide.slide_id}: Removed port 8000")
                    print(f"  [OK] Removed 'port 8000'")

                # A10: Add "9 AI Agents" naming
                if "7 agents" in text.lower() or "7 AI" in text:
                    new_text = text.replace("7 agents", "9 agents").replace("7 AI", "9 AI")
                    shape.text = new_text
                    changes.append(f"Slide {slide.slide_id}: 7 => 9 agents")
                    print(f"  [OK] Updated agent count to 9")

                # A11: Update stats
                if "8 docs" in text.lower() or "8 documents" in text.lower():
                    new_text = text.replace("8 docs", "30+ outputs").replace("8 documents", "30+ outputs")
                    shape.text = new_text
                    changes.append(f"Slide {slide.slide_id}: 8 docs => 30+ outputs")
                    print(f"  [OK] Updated output count")

                # A13: Replace "Swagger · CORS" with "OpenAPI docs · Secure API"
                if "CORS" in text:
                    new_text = text.replace("Swagger docs · CORS", "OpenAPI docs · Secure CORS config")
                    new_text = new_text.replace("Swagger · CORS", "OpenAPI docs · Secure CORS config")
                    shape.text = new_text
                    changes.append(f"Slide {slide.slide_id}: Updated Swagger/CORS")
                    print(f"  [OK] Updated Swagger/CORS reference")

    # === SLIDE 5: Expected Outcomes ===
    print("\n=== SLIDE 5: Outcomes ===")

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text

                # A16: Update "$0 infra"
                if "$0 Infrastructure" in text or "$0 infra" in text:
                    new_text = text.replace("$0 Infrastructure cost", "INR0 cloud hosting / ~INR500 per project via API")
                    new_text = new_text.replace("$0 infra", "INR0 cloud hosting")
                    shape.text = new_text
                    changes.append(f"Slide {slide.slide_id}: Updated infrastructure cost")
                    print(f"  [OK] Updated infrastructure cost description")

                # Update NOW Delivered section with P7a and P8c full scope
                if "NOW Delivered" in text or "8-phase pipeline" in text:
                    # Update to reflect 11 phases and 9 agents
                    new_text = text.replace("8-phase pipeline", "11-phase pipeline")
                    new_text = new_text.replace("8-phase", "11-phase")
                    shape.text = new_text
                    changes.append(f"Slide {slide.slide_id}: Updated phase count")
                    print(f"  [OK] Updated to 11-phase pipeline")

    # === MISRA-C dual-version ===
    print("\n=== MISRA-C ===")

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                if "MISRA-C" in text and ":2012 / 2023" not in text:
                    new_text = text.replace("MISRA-C", "MISRA-C:2012 / 2023")
                    shape.text = new_text
                    changes.append(f"Slide {slide.slide_id}: Added MISRA-C version detail")
                    print(f"  [OK] Updated MISRA-C to dual-version")

    # === P8 label update ===
    print("\n=== P8 Label ===")

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                # Update P8 descriptions to show full scope
                if "P8 -" in text or "P8:" in text:
                    if "Code Review" in text:
                        new_text = text.replace("SRS + SDD + Code Review",
                                                          "Code Gen + Qt GUI + CI/CD => Git PR")
                        shape.text = new_text
                        changes.append(f"Slide {slide.slide_id}: Updated P8 scope description")
                        print(f"  [OK] Updated P8 to show full scope")

    # === Summary ===
    print(f"\n=== Summary ===")
    print(f"Total changes made: {len(changes)}")
    print("\nChanges:")
    for i, change in enumerate(changes, 1):
        print(f"  {i}. {change.encode('ascii', 'replace').decode('ascii')}")

    # Save the updated presentation
    output_path = pptx_path.parent / f"{pptx_path.stem}_updated{pptx_path.suffix}"
    print(f"\nSaving updated file to: {output_path}")
    prs.save(str(output_path))
    print(f"[OK] Done! File saved.")

    return output_path

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python update_pptx.py <path_to_pptx>")
        print("\nExample:")
        print("  python update_pptx.py HardwarePipeline_FINAL_v3.pptx")
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    if not pptx_path.exists():
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)

    update_pptx(pptx_path)
