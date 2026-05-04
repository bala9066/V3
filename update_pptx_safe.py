"""
Safe PPTX updater - preserves all content while making targeted text replacements
"""

import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Installing python-pptx...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation

def safe_update_pptx(pptx_path: Path):
    """Update PPTX with minimal changes to preserve content"""

    print(f"Loading {pptx_path}...")
    prs = Presentation(str(pptx_path))
    print(f"Slides: {len(prs.slides)}")
    print(f"Slide size: {prs.slide_width}x{prs.slide_height}")

    changes = []

    # Define all replacements (old_text -> new_text)
    replacements = [
        # Phase count
        ("8 Phases", "11 Phases"),
        ("8 phases", "11 phases"),
        ("8-phase", "11-phase"),

        # Agent count
        ("7 agents", "9 agents"),
        ("7 AI", "9 AI"),
        ("7-agent", "9-agent"),

        # Output count
        ("8 docs", "30+ outputs"),
        ("8 documents", "30+ outputs"),
        ("8-docs", "30+ outputs"),

        # Standards
        ("FCC, MISRA-C", "MIL-STD-461/810, MISRA-C"),
        ("IEEE 29148, IEEE 830, IEEE 1016", "IEEE 29148, IEEE 830, IEEE 1016, MISRA-C"),

        # Problem phrasing
        ("from requirements to code review", "from requirements to production-ready firmware"),

        # Infrastructure
        ("$0 Infrastructure cost", "INR 0 cloud hosting"),
        ("$0 infra", "INR 0 cloud hosting"),

        # P8 scope
        ("SRS + SDD + Code Review", "Code Gen + Qt GUI + CI/CD -> Git PR"),
        ("SRS+SDD+Code Review", "Code Gen + Qt GUI + CI/CD -> Git PR"),

        # Remove port
        (" port 8000", ""),
        (": 8000", ""),

        # Swagger/CORS
        ("Swagger · CORS", "OpenAPI docs · Secure CORS config"),
        ("Swagger docs · CORS", "OpenAPI docs · Secure CORS config"),

        # MISRA-C
        ("MISRA-C:", "MISRA-C:2012 / 2023"),
    ]

    # Process each slide
    for slide_idx, slide in enumerate(prs.slides, 1):
        print(f"\nSlide {slide_idx}:")

        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue

            for paragraph in shape.text_frame.paragraphs:
                original_text = paragraph.text

                if not original_text:
                    continue

                modified_text = original_text

                # Apply all replacements
                for old, new in replacements:
                    if old in modified_text:
                        modified_text = modified_text.replace(old, new)

                # Only update if changed
                if modified_text != original_text:
                    paragraph.text = modified_text
                    changes.append(f"Slide {slide_idx}: '{original_text[:40]}...' -> '{modified_text[:40]}...'")
                    print(f"  [OK] Updated: {original_text[:40]}...")

    print(f"\n\n{'='*60}")
    print(f"Total changes: {len(changes)}")
    print(f"{'='*60}")

    # Save with same name (overwrite)
    print(f"\nSaving to: {pptx_path}")
    prs.save(str(pptx_path))
    print("[OK] PPTX updated successfully!")

    return len(changes)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python update_pptx_safe.py <path_to_pptx>")
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    if not pptx_path.exists():
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)

    safe_update_pptx(pptx_path)
